# backend/utils/extractors.py
import fitz  # PyMuPDF for PDFs
import docx
from pptx import Presentation

def extract_pdf(file_path):
    doc = fitz.open(file_path)
    return " ".join([page.get_text() for page in doc])

def extract_docx(file_path):
    doc = docx.Document(file_path)
    return " ".join([p.text for p in doc.paragraphs])

def extract_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return " ".join(text_runs)
